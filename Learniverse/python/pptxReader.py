import sys
from pptx import Presentation

def extract_text_from_pptx(pptx_file: str) -> str:
    try:
        prs = Presentation(pptx_file)
        extracted_text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text.append(shape.text)
        
        return '\n'.join(extracted_text)

    except Exception as e:
        print(str(e), file=sys.stderr)
        return ''

if __name__ == '__main__':
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        pptx_text = extract_text_from_pptx(file_path)
        print(pptx_text)
    else:
        print("Please provide a PPTX file path.")
